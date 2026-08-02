"""Text extraction for every upload type the RAG module accepts.

PDF / Word / Excel / PPT / plain text are handled with real parsers.
Images go through RapidOCR (ONNX Runtime) — see app/modules/vision/ocr.py
for why that's used instead of Tesseract on this host.
"""

import io
from pathlib import Path

import pypdf
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation


class UnsupportedFileType(ValueError):
    pass


def _extract_pdf(content: bytes) -> str:
    reader = pypdf.PdfReader(io.BytesIO(content))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_docx(content: bytes) -> str:
    doc = DocxDocument(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_xlsx(content: bytes) -> str:
    wb = load_workbook(io.BytesIO(content), data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        parts.append(f"# Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(content: bytes) -> str:
    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _extract_image(content: bytes) -> str:
    from app.modules.vision.ocr import extract_text as ocr_extract_text

    return ocr_extract_text(content)


_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".xlsm": _extract_xlsx,
    ".pptx": _extract_pptx,
    ".png": _extract_image,
    ".jpg": _extract_image,
    ".jpeg": _extract_image,
}

_TEXT_EXTENSIONS = {".txt", ".md", ".csv", ".json", ".log"}


def extract_text(filename: str, content: bytes) -> str:
    ext = Path(filename).suffix.lower()

    if ext in _TEXT_EXTENSIONS:
        return content.decode("utf-8", errors="replace")

    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        raise UnsupportedFileType(
            f"Unsupported file type '{ext}'. Supported: "
            f"{sorted(_EXTRACTORS) + sorted(_TEXT_EXTENSIONS)}"
        )
    return extractor(content)


def extract_text_from_url(url: str) -> str:
    import requests
    from bs4 import BeautifulSoup

    resp = requests.get(url, timeout=15, headers={"User-Agent": "OmniAI-RAG/0.1"})
    resp.raise_for_status()

    soup = BeautifulSoup(resp.text, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()

    text = soup.get_text(separator="\n")
    lines = [line.strip() for line in text.splitlines()]
    return "\n".join(line for line in lines if line)
